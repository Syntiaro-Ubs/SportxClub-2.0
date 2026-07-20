"""
SportXClub 2.0 — Payment Process Architecture
Professional PowerPoint Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# ── BRAND COLORS ──
BRAND_GREEN = RGBColor(0x6D, 0xFF, 0x3B)
DARK_BG = RGBColor(0x05, 0x05, 0x05)
CARD_BG = RGBColor(0x10, 0x12, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_DIM = RGBColor(0x94, 0xA3, 0xB8)
WHITE_MID = RGBColor(0xCB, 0xD5, 0xE1)
EMERALD = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
LIGHT_GREEN = RGBColor(0x34, 0xD3, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1), corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or CARD_BG
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_label(slide, left, top, text, font_size=9, color=BRAND_GREEN):
    return add_text_box(slide, left, top, Inches(4), Inches(0.3), text, font_size=font_size, color=color, bold=True)

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=WHITE_MID, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
    return txBox

# ════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1)

# Green accent circle (top right)
circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5), Inches(4), Inches(4))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x0A)
circle.line.fill.background()

# Brand dot
dot = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(1.5), Inches(0.15), Inches(0.15))
dot.fill.solid()
dot.fill.fore_color.rgb = BRAND_GREEN
dot.line.fill.background()

add_text_box(slide1, Inches(1.3), Inches(1.4), Inches(3), Inches(0.4), "SPORTXCLUB", font_size=14, color=WHITE, bold=True)

add_label(slide1, Inches(1.0), Inches(2.5), "PAYMENT PROCESS ARCHITECTURE  ·  JULY 2026", font_size=10)

add_text_box(slide1, Inches(1.0), Inches(3.0), Inches(8), Inches(1.5),
    "Payment Workflow &\nSettlement Engine", font_size=44, color=WHITE, bold=True)

add_text_box(slide1, Inches(1.0), Inches(4.8), Inches(7.5), Inches(1),
    "Complete end-to-end payment architecture covering Player payments,\nPlatform commission model, Owner settlements, and refund workflows\n— powered by Razorpay Route on Node.js.",
    font_size=14, color=WHITE_DIM)

# Tags
tag1 = add_shape(slide1, Inches(1.0), Inches(6.2), Inches(2.8), Inches(0.4), fill_color=BRAND_GREEN)
add_text_box(slide1, Inches(1.1), Inches(6.22), Inches(2.6), Inches(0.36), "Presented by Team SportXClub", font_size=11, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)

tag2 = add_shape(slide1, Inches(4.0), Inches(6.2), Inches(2.5), Inches(0.4), fill_color=RGBColor(0x15, 0x18, 0x1E), border_color=RGBColor(0x30, 0x33, 0x3A))
add_text_box(slide1, Inches(4.1), Inches(6.22), Inches(2.3), Inches(0.36), "Confidential · Internal Use", font_size=11, color=WHITE_DIM, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 2: CURRENT STATE vs PROPOSED
# ════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2)

add_label(slide2, Inches(0.8), Inches(0.5), "CURRENT STATE ANALYSIS")
add_text_box(slide2, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "What We Have vs What We Need", font_size=32, bold=True)

# Left column - Current (Red)
add_text_box(slide2, Inches(0.8), Inches(1.8), Inches(5), Inches(0.4), "❌  Current State (Mock Only)", font_size=15, color=ROSE, bold=True)

current_items = [
    "Payment simulated with setTimeout(2000)",
    "No real payment gateway integration",
    "Hardcoded price breakdown (₹45 fee, ₹135 GST)",
    "No owner settlement or commission tracking",
    "Split payment is UI-only (no real logic)",
    "No webhook verification or transaction records"
]
for i, item in enumerate(current_items):
    y = Inches(2.35) + Inches(i * 0.55)
    card = add_shape(slide2, Inches(0.8), y, Inches(5.4), Inches(0.45), fill_color=RGBColor(0x14, 0x0A, 0x0E), border_color=RGBColor(0x3D, 0x14, 0x1C))
    add_text_box(slide2, Inches(1.0), y + Inches(0.06), Inches(5), Inches(0.35), f"•  {item}", font_size=12, color=RGBColor(0xFB, 0x71, 0x85))

# Right column - Proposed (Green)
add_text_box(slide2, Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.4), "✅  Proposed State (Production)", font_size=15, color=BRAND_GREEN, bold=True)

proposed_items = [
    "Razorpay Route with auto-split to owners",
    "Dynamic pricing (base + GST + fee − coupon)",
    "Automated T+2 settlement via cron jobs",
    "15% platform commission (configurable tiers)",
    "Real split payments with 30-min timeout",
    "Webhook verification + full audit trail"
]
for i, item in enumerate(proposed_items):
    y = Inches(2.35) + Inches(i * 0.55)
    card = add_shape(slide2, Inches(6.8), y, Inches(5.7), Inches(0.45), fill_color=RGBColor(0x06, 0x14, 0x0A), border_color=RGBColor(0x14, 0x3D, 0x1C))
    add_text_box(slide2, Inches(7.0), y + Inches(0.06), Inches(5.3), Inches(0.35), f"•  {item}", font_size=12, color=LIGHT_GREEN)


# ════════════════════════════════════════════════════
# SLIDE 3: THREE STAKEHOLDERS
# ════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3)

add_label(slide3, Inches(0.8), Inches(0.5), "STAKEHOLDER OVERVIEW")
add_text_box(slide3, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Three Parties, One Flow", font_size=32, bold=True)
add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(9), Inches(0.4),
    "Every booking involves three stakeholders with distinct roles in the payment lifecycle.",
    font_size=14, color=WHITE_DIM)

stakeholders = [
    {"emoji": "👤", "title": "Player (User)", "desc": "Browses venues, selects time slots,\nmakes payment via UPI/Card/Wallet.\nCan split payment with teammates.", "tag": "PAYS MONEY →", "tag_color": BLUE, "x": 0.8},
    {"emoji": "🏢", "title": "SportXClub (Platform)", "desc": "Middleman. Collects payment,\ndeducts 15% commission + ₹45 fee,\nsettles remainder to owner.", "tag": "HOLDS & SPLITS", "tag_color": BRAND_GREEN, "x": 4.8},
    {"emoji": "🏟️", "title": "Turf Owner", "desc": "Receives 85% of base price after\ncommission deduction. Settlement\nvia Razorpay Route in T+2 days.", "tag": "← RECEIVES MONEY", "tag_color": VIOLET, "x": 8.8},
]

for s in stakeholders:
    x = Inches(s["x"])
    card = add_shape(slide3, x, Inches(2.3), Inches(3.6), Inches(4.2), border_color=RGBColor(0x25, 0x28, 0x2E))
    add_text_box(slide3, x + Inches(0.3), Inches(2.6), Inches(3), Inches(0.6), s["emoji"], font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x + Inches(0.3), Inches(3.3), Inches(3), Inches(0.4), s["title"], font_size=18, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x + Inches(0.3), Inches(3.9), Inches(3), Inches(1.2), s["desc"], font_size=12, color=WHITE_DIM, alignment=PP_ALIGN.CENTER)
    tag_shape = add_shape(slide3, x + Inches(0.5), Inches(5.5), Inches(2.6), Inches(0.4), fill_color=RGBColor(0x0A, 0x12, 0x0A), border_color=s["tag_color"])
    add_text_box(slide3, x + Inches(0.5), Inches(5.52), Inches(2.6), Inches(0.36), s["tag"], font_size=10, color=s["tag_color"], bold=True, alignment=PP_ALIGN.CENTER)

# Arrows between cards
add_text_box(slide3, Inches(4.35), Inches(3.8), Inches(0.5), Inches(0.5), "→", font_size=28, color=BRAND_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide3, Inches(8.35), Inches(3.8), Inches(0.5), Inches(0.5), "→", font_size=28, color=BRAND_GREEN, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 4: MONEY FLOW & BREAKDOWN
# ════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4)

add_label(slide4, Inches(0.8), Inches(0.5), "REVENUE MODEL")
add_text_box(slide4, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Where Does The Money Go?", font_size=32, bold=True)
add_text_box(slide4, Inches(0.8), Inches(1.6), Inches(6), Inches(0.4), "Price breakdown for a ₹1,000 base booking:", font_size=14, color=WHITE_DIM)

# Table data
table_data = [
    ["Component", "Amount", "Recipient"],
    ["Base Turf Price", "₹1,000", "Split below ↓"],
    ["→ Owner Share (85%)", "₹850", "Turf Owner"],
    ["→ Platform Commission (15%)", "₹150", "SportXClub"],
    ["Convenience Fee", "₹45", "SportXClub"],
    ["GST (18%)", "₹188", "Government"],
    ["Coupon Discount", "−₹100", "SportXClub absorbs"],
    ["TOTAL USER PAYS", "₹1,133", "—"],
]

rows, cols = len(table_data), 3
tbl = slide4.shapes.add_table(rows, cols, Inches(0.8), Inches(2.1), Inches(6), Inches(3.8)).table

tbl.columns[0].width = Inches(2.8)
tbl.columns[1].width = Inches(1.4)
tbl.columns[2].width = Inches(1.8)

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        cell.fill.solid()

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.name = 'Calibri'

        if r == 0:
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x18, 0x1E)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = BRAND_GREEN
                p.font.bold = True
                p.font.size = Pt(10)
        elif r == len(table_data) - 1:
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x0A)
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.color.rgb = BRAND_GREEN
        else:
            cell.fill.fore_color.rgb = RGBColor(0x0C, 0x0E, 0x12)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE_MID
                if c == 1 and "₹850" in table_data[r][c]:
                    p.font.color.rgb = LIGHT_GREEN
                    p.font.bold = True
                elif c == 1 and "₹150" in table_data[r][c]:
                    p.font.color.rgb = BRAND_GREEN
                    p.font.bold = True
                elif c == 1 and "−" in table_data[r][c]:
                    p.font.color.rgb = ROSE

# Pie chart on the right
chart_data = CategoryChartData()
chart_data.categories = ['Owner (54%)', 'Commission (15%)', 'GST (15%)', 'Coupon (10%)', 'Gateway (6%)']
chart_data.add_series('Distribution', (54, 15, 15, 10, 6))

chart_frame = slide4.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(7.5), Inches(1.8), Inches(5), Inches(5),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = WHITE_DIM

plot = chart.plots[0]
series = plot.series[0]
colors = [BRAND_GREEN, EMERALD, AMBER, ROSE, BLUE]
for i, color in enumerate(colors):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color

chart_frame.chart.chart_style = 2


# ════════════════════════════════════════════════════
# SLIDE 5: PLAYER PAYMENT JOURNEY
# ════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5)

add_label(slide5, Inches(0.8), Inches(0.5), "PLAYER SIDE")
add_text_box(slide5, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Player Booking → Payment Journey", font_size=32, bold=True)

steps = [
    ("Step 1", "Browse\nVenues"),
    ("Step 2", "Select Sport\n& Slot"),
    ("Step 3", "Review\nPricing"),
    ("Step 4", "Choose\nPayment"),
    ("Step 5", "Razorpay\nCheckout"),
    ("Done ✅", "Booking\nConfirmed"),
]

for i, (num, label) in enumerate(steps):
    x = Inches(0.6 + i * 2.1)
    is_highlight = i == 3
    is_done = i == 5
    bg = RGBColor(0x0A, 0x1A, 0x0A) if is_highlight else (RGBColor(0x06, 0x14, 0x12) if is_done else CARD_BG)
    border = BRAND_GREEN if is_highlight else (EMERALD if is_done else RGBColor(0x25, 0x28, 0x2E))

    card = add_shape(slide5, x, Inches(2.0), Inches(1.8), Inches(1.3), fill_color=bg, border_color=border)
    add_text_box(slide5, x, Inches(2.1), Inches(1.8), Inches(0.3), num, font_size=9, color=BRAND_GREEN if not is_done else EMERALD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, x, Inches(2.45), Inches(1.8), Inches(0.7), label, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    if i < 5:
        add_text_box(slide5, x + Inches(1.75), Inches(2.35), Inches(0.4), Inches(0.5), "→", font_size=22, color=BRAND_GREEN, alignment=PP_ALIGN.CENTER)

# Stats cards
stats = [
    ("3", "Payment Methods\n(UPI, Card, Wallet)", BRAND_GREEN),
    ("30 min", "Split Payment\nTimeout Window", AMBER),
    ("QR", "Entry Pass with\nQR Code Generated", EMERALD),
    ("PDF", "Downloadable\nReceipt with GST", BLUE),
]

for i, (val, desc, clr) in enumerate(stats):
    x = Inches(0.6 + i * 3.1)
    card = add_shape(slide5, x, Inches(4.0), Inches(2.8), Inches(1.5), border_color=RGBColor(0x25, 0x28, 0x2E))
    add_text_box(slide5, x + Inches(0.2), Inches(4.15), Inches(2.4), Inches(0.5), val, font_size=26, color=clr, bold=True)
    add_text_box(slide5, x + Inches(0.2), Inches(4.7), Inches(2.4), Inches(0.7), desc, font_size=11, color=WHITE_DIM)

# Payment states
add_text_box(slide5, Inches(0.8), Inches(5.8), Inches(12), Inches(0.3),
    "Payment States:   INITIATED  →  PROCESSING  →  SUCCESS ✅  |  FAILED ❌  |  PARTIALLY_PAID ⏳  |  REFUNDED 🔄",
    font_size=11, color=WHITE_DIM)


# ════════════════════════════════════════════════════
# SLIDE 6: INTERNAL TECH ARCHITECTURE
# ════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6)

add_label(slide6, Inches(0.8), Inches(0.5), "TECHNICAL ARCHITECTURE")
add_text_box(slide6, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Internal Payment Engine — Node.js", font_size=32, bold=True)

# Left: Timeline
timeline_items = [
    ("1. Frontend → Create Order API", "React app calls /api/payment/create-order with booking details"),
    ("2. Node.js → Razorpay Order + Route", "Server creates Order with transfers[] array specifying owner's linked account"),
    ("3. Frontend → Razorpay Checkout", "Razorpay Checkout.js popup opens — player pays via UPI/Card"),
    ("4. Server → Verify Signature", "HMAC-SHA256 signature verification to confirm payment is genuine"),
    ("5. Razorpay → Webhook Notification", "Backup confirmation via payment.captured webhook event"),
    ("6. Cron Job → Release & Settle", "Nightly batch releases held funds to owner's bank (T+2)"),
]

for i, (title, desc) in enumerate(timeline_items):
    y = Inches(1.85 + i * 0.85)
    dot = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y + Inches(0.08), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BRAND_GREEN
    dot.line.fill.background()

    if i < len(timeline_items) - 1:
        line_shape = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y + Inches(0.2), Inches(0.02), Inches(0.7))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x3D, 0x1A)
        line_shape.line.fill.background()

    add_text_box(slide6, Inches(1.3), y - Inches(0.02), Inches(5), Inches(0.3), title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide6, Inches(1.3), y + Inches(0.28), Inches(5.2), Inches(0.4), desc, font_size=11, color=WHITE_DIM)

# Right: Tech stack
stack_items = [
    ("⚡", "Node.js + Express.js", "API Server & Business Logic", BRAND_GREEN),
    ("💳", "Razorpay Route API", "Payment Gateway + Auto-Split", BLUE),
    ("🗃️", "MongoDB", "Transactions, Settlements, Bookings", VIOLET),
    ("⏰", "node-cron", "Nightly Settlement Batch Jobs", AMBER),
    ("🔔", "Webhooks + Notifications", "Email (Nodemailer) + SMS (MSG91)", ROSE),
]

add_text_box(slide6, Inches(7.2), Inches(1.7), Inches(5), Inches(0.3), "Tech Stack", font_size=14, color=WHITE_DIM, bold=True)

for i, (emoji, name, desc, clr) in enumerate(stack_items):
    y = Inches(2.15 + i * 0.95)
    card = add_shape(slide6, Inches(7.2), y, Inches(5.2), Inches(0.8), border_color=RGBColor(0x25, 0x28, 0x2E))

    icon_bg = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.45), y + Inches(0.15), Inches(0.5), Inches(0.5))
    icon_bg.fill.solid()
    icon_bg.fill.fore_color.rgb = RGBColor(0x15, 0x18, 0x1E)
    icon_bg.line.fill.background()
    add_text_box(slide6, Inches(7.45), y + Inches(0.12), Inches(0.5), Inches(0.5), emoji, font_size=16, alignment=PP_ALIGN.CENTER)

    add_text_box(slide6, Inches(8.15), y + Inches(0.1), Inches(4), Inches(0.3), name, font_size=13, color=WHITE, bold=True)
    add_text_box(slide6, Inches(8.15), y + Inches(0.4), Inches(4), Inches(0.3), desc, font_size=10, color=WHITE_DIM)


# ════════════════════════════════════════════════════
# SLIDE 7: OWNER SETTLEMENT
# ════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7)

add_label(slide7, Inches(0.8), Inches(0.5), "OWNER SIDE")
add_text_box(slide7, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Owner Settlement Lifecycle", font_size=32, bold=True)

# Flow steps
settle_steps = [
    ("Booking", "Player\nPays"),
    ("Hold", "Platform\nEscrow"),
    ("Deduct", "15%\nCommission"),
    ("T+2", "Bank\nTransfer"),
    ("Done ✅", "Owner\nReceives"),
]

for i, (num, label) in enumerate(settle_steps):
    x = Inches(0.5 + i * 2.5)
    is_done = i == 4
    is_hold = i == 1
    bg = RGBColor(0x0A, 0x1A, 0x0A) if is_hold else (RGBColor(0x06, 0x14, 0x12) if is_done else CARD_BG)
    border = BRAND_GREEN if is_hold else (EMERALD if is_done else RGBColor(0x25, 0x28, 0x2E))

    card = add_shape(slide7, x, Inches(1.9), Inches(2.0), Inches(1.2), fill_color=bg, border_color=border)
    add_text_box(slide7, x, Inches(2.0), Inches(2.0), Inches(0.25), num, font_size=9, color=BRAND_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, x, Inches(2.3), Inches(2.0), Inches(0.6), label, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    if i < 4:
        add_text_box(slide7, x + Inches(1.95), Inches(2.2), Inches(0.5), Inches(0.5), "→", font_size=22, color=BRAND_GREEN, alignment=PP_ALIGN.CENTER)

# Settlement table
settle_tbl_data = [
    ["Booking", "Player", "Gross", "Commission (15%)", "Net Payout", "Status"],
    ["#SX-001", "Rahul Sharma", "₹2,400", "₹360", "₹2,040", "✅ Settled"],
    ["#SX-002", "Priya Patel", "₹1,500", "₹225", "₹1,275", "✅ Settled"],
    ["#SX-003", "Arjun M.", "₹800", "₹120", "₹680", "⏳ Processing"],
    ["#SX-004", "Vikram S.", "₹1,200", "—", "₹0", "❌ Cancelled"],
]

rows, cols = len(settle_tbl_data), 6
tbl = slide7.shapes.add_table(rows, cols, Inches(0.5), Inches(3.5), Inches(12), Inches(3)).table

for c in range(cols):
    tbl.columns[c].width = Inches(2)

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = settle_tbl_data[r][c]
        cell.fill.solid()
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = 'Calibri'
        if r == 0:
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x18, 0x1E)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = BRAND_GREEN
                p.font.bold = True
                p.font.size = Pt(10)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x0C, 0x0E, 0x12)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE_MID
            if c == 5:
                for p in cell.text_frame.paragraphs:
                    if "Settled" in settle_tbl_data[r][c]:
                        p.font.color.rgb = LIGHT_GREEN
                    elif "Processing" in settle_tbl_data[r][c]:
                        p.font.color.rgb = AMBER
                    elif "Cancelled" in settle_tbl_data[r][c]:
                        p.font.color.rgb = ROSE
                    p.font.bold = True


# ════════════════════════════════════════════════════
# SLIDE 8: REFUND POLICY
# ════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8)

add_label(slide8, Inches(0.8), Inches(0.5), "CANCELLATION & REFUND")
add_text_box(slide8, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Refund Policy — 3 Tier Model", font_size=32, bold=True)

refund_tiers = [
    {"pct": "100%", "title": "Free Cancellation", "desc": "Cancelled more than 4 hours\nbefore slot time.\nFull refund in 5-7 business days.", "sub": "Owner: ₹0 · Platform: ₹0", "color": EMERALD, "emoji": "✅"},
    {"pct": "50%", "title": "Late Cancellation", "desc": "Cancelled 1-4 hours before slot.\nHalf amount refunded,\nhalf distributed.", "sub": "Owner: 50% · Platform: 15% of 50%", "color": AMBER, "emoji": "⚠️"},
    {"pct": "0%", "title": "No-Show / Late", "desc": "Cancelled within 1 hour or\ndidn't show up.\nNo refund issued.", "sub": "Owner: 100% · Platform: Full 15%", "color": ROSE, "emoji": "🚫"},
]

for i, tier in enumerate(refund_tiers):
    x = Inches(0.6 + i * 4.1)
    bg_color = RGBColor(0x06, 0x14, 0x0A) if i == 0 else (RGBColor(0x14, 0x10, 0x06) if i == 1 else RGBColor(0x14, 0x0A, 0x0E))
    card = add_shape(slide8, x, Inches(1.9), Inches(3.8), Inches(4.5), fill_color=bg_color, border_color=tier["color"])

    add_text_box(slide8, x, Inches(2.15), Inches(3.8), Inches(0.5), tier["emoji"], font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, x, Inches(2.7), Inches(3.8), Inches(0.6), tier["pct"], font_size=48, color=tier["color"], bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, x, Inches(3.4), Inches(3.8), Inches(0.4), tier["title"], font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, x + Inches(0.3), Inches(3.9), Inches(3.2), Inches(1), tier["desc"], font_size=12, color=WHITE_DIM, alignment=PP_ALIGN.CENTER)
    
    sub_shape = add_shape(slide8, x + Inches(0.3), Inches(5.3), Inches(3.2), Inches(0.4), fill_color=RGBColor(0x0A, 0x0C, 0x10), border_color=RGBColor(0x20, 0x23, 0x28))
    add_text_box(slide8, x + Inches(0.3), Inches(5.32), Inches(3.2), Inches(0.36), tier["sub"], font_size=10, color=tier["color"], bold=True, alignment=PP_ALIGN.CENTER)

# Timeline
add_shape(slide8, Inches(0.6), Inches(6.1), Inches(12), Inches(0.6), border_color=RGBColor(0x25, 0x28, 0x2E))
add_text_box(slide8, Inches(0.8), Inches(6.15), Inches(0.3), Inches(0.5), "⏱️", font_size=18)
add_text_box(slide8, Inches(1.2), Inches(6.18), Inches(10), Inches(0.4),
    "Refund Timeline:  Day 0: Initiated  →  Day 1: Razorpay processes  →  Day 3-5: Credited to bank  →  Day 5-7: Worst case",
    font_size=11, color=WHITE_DIM)


# ════════════════════════════════════════════════════
# SLIDE 9: REVENUE PROJECTION CHART
# ════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9)

add_label(slide9, Inches(0.8), Inches(0.5), "FINANCIAL PROJECTION")
add_text_box(slide9, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Monthly Revenue Forecast", font_size=32, bold=True)
add_text_box(slide9, Inches(0.8), Inches(1.55), Inches(8), Inches(0.4), "Based on 100 bookings/day at ₹1,500 average base price.", font_size=14, color=WHITE_DIM)

# Bar chart
chart_data2 = CategoryChartData()
chart_data2.categories = ['Month 1', 'Month 2', 'Month 3', 'Month 4', 'Month 5', 'Month 6']
chart_data2.add_series('Platform Revenue (₹ Lakhs)', (2.7, 4.5, 6.1, 7.2, 8.1, 9.5))

chart_frame2 = slide9.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2.1), Inches(7.5), Inches(4.8),
    chart_data2
)
chart2 = chart_frame2.chart
chart2.has_legend = False
chart2.chart_style = 2

plot2 = chart2.plots[0]
series2 = plot2.series[0]
series2.format.fill.solid()
series2.format.fill.fore_color.rgb = BRAND_GREEN

chart2.category_axis.tick_labels.font.color.rgb = WHITE_DIM
chart2.category_axis.tick_labels.font.size = Pt(10)
chart2.value_axis.tick_labels.font.color.rgb = WHITE_DIM
chart2.value_axis.tick_labels.font.size = Pt(10)

# Right side stats
rev_stats = [
    ("📊", "Monthly GMV (at scale)", "₹45L", WHITE),
    ("💰", "Platform Revenue", "₹8.1L", BRAND_GREEN),
    ("💳", "Razorpay Costs", "₹87K", ROSE),
    ("🚀", "Net Profit / Month", "₹7.2L", LIGHT_GREEN),
]

for i, (emoji, label, value, clr) in enumerate(rev_stats):
    y = Inches(2.1 + i * 1.2)
    card = add_shape(slide9, Inches(8.5), y, Inches(4.2), Inches(1.0), border_color=RGBColor(0x25, 0x28, 0x2E))
    add_text_box(slide9, Inches(8.7), y + Inches(0.1), Inches(0.5), Inches(0.4), emoji, font_size=20)
    add_text_box(slide9, Inches(9.3), y + Inches(0.08), Inches(3), Inches(0.3), label, font_size=11, color=WHITE_DIM)
    add_text_box(slide9, Inches(9.3), y + Inches(0.42), Inches(3), Inches(0.4), value, font_size=26, color=clr, bold=True)


# ════════════════════════════════════════════════════
# SLIDE 10: DOCUMENTS & NEXT STEPS
# ════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10)

add_label(slide10, Inches(0.8), Inches(0.5), "IMPLEMENTATION ROADMAP")
add_text_box(slide10, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7), "Documents Required & Next Steps", font_size=32, bold=True)

# Left: Documents
add_text_box(slide10, Inches(0.8), Inches(1.8), Inches(5), Inches(0.3), "📋  Documents Checklist", font_size=14, color=BRAND_GREEN, bold=True)

doc_items = [
    ("🏢", "Platform (SportXClub)", "Business Registration · PAN · GST ·\nBank Account · Razorpay Route Activation"),
    ("🏟️", "Per Turf Owner", "PAN · Aadhaar · Bank Details ·\nCancelled Cheque · Turf Ownership Proof"),
    ("👤", "Player (User)", "Name · Email · Phone only\n— no documents required"),
    ("⚖️", "Legal Pages", "Privacy Policy · Terms of Service ·\nRefund Policy · Contact Page"),
]

for i, (emoji, title, desc) in enumerate(doc_items):
    y = Inches(2.2 + i * 1.15)
    card = add_shape(slide10, Inches(0.8), y, Inches(5.5), Inches(1.0), border_color=RGBColor(0x25, 0x28, 0x2E))
    add_text_box(slide10, Inches(1.0), y + Inches(0.08), Inches(0.5), Inches(0.4), emoji, font_size=18)
    add_text_box(slide10, Inches(1.5), y + Inches(0.08), Inches(4.5), Inches(0.3), title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide10, Inches(1.5), y + Inches(0.38), Inches(4.5), Inches(0.55), desc, font_size=10, color=WHITE_DIM)

# Right: Timeline
add_text_box(slide10, Inches(7.0), Inches(1.8), Inches(5), Inches(0.3), "🚀  5-Week Rollout Plan", font_size=14, color=BRAND_GREEN, bold=True)

weeks = [
    ("Week 1-2", "Business Setup", "LLP registration, GST, bank account, CA"),
    ("Week 2-3", "Razorpay Integration", "Account, Route activation, test mode"),
    ("Week 3-4", "Backend Development", "Node.js APIs, webhooks, cron jobs"),
    ("Week 4", "Owner Onboarding", "KYC flow, linked accounts, admin panel"),
    ("Week 5", "Go Live! 🎉", "Switch to live, first real transactions"),
]

for i, (week, title, desc) in enumerate(weeks):
    y = Inches(2.2 + i * 0.9)

    dot = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.1), y + Inches(0.12), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BRAND_GREEN
    dot.line.fill.background()

    if i < len(weeks) - 1:
        line = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.15), y + Inches(0.24), Inches(0.02), Inches(0.72))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x1A, 0x3D, 0x1A)
        line.line.fill.background()

    add_text_box(slide10, Inches(7.4), y, Inches(1.2), Inches(0.3), week, font_size=10, color=BRAND_GREEN, bold=True)
    add_text_box(slide10, Inches(7.4), y + Inches(0.25), Inches(4.5), Inches(0.3), title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide10, Inches(7.4), y + Inches(0.52), Inches(4.5), Inches(0.3), desc, font_size=10, color=WHITE_DIM)

# Cost box
cost_card = add_shape(slide10, Inches(7.0), Inches(6.2), Inches(5.5), Inches(0.7), fill_color=RGBColor(0x0A, 0x1A, 0x0A), border_color=BRAND_GREEN)
add_text_box(slide10, Inches(7.2), Inches(6.25), Inches(5), Inches(0.3), "💰  Startup Cost:  ₹10,000 — ₹15,000", font_size=14, color=BRAND_GREEN, bold=True)
add_text_box(slide10, Inches(7.2), Inches(6.55), Inches(5), Inches(0.25), "Razorpay is FREE. Major cost = LLP registration + CA fees.", font_size=10, color=WHITE_DIM)


# ════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════
output_path = r"d:\VS_CODE_SPRING\SportXClub 2.0\SportxClub-2.0\SportXClub-Payment-Presentation.pptx"
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
